#!/usr/bin/env python3
"""pptx_to_chunks.py — extract slides from a PPTX deck as text chunks.

Each slide → one chunk: title from layout placeholder + body text. Notes
(speaker notes) are attached as a separate field. Useful for ingesting course
slide decks / industry-conf presentations.

Lazy install: python-pptx.

Usage:
  python3 pptx_to_chunks.py <input.pptx> [output.jsonl|-]

Output JSONL (one slide per line):
  {"source": "...", "slide_index": 1, "title": "...", "body": "...",
   "notes": "...", "word_count": ...}

Exit codes: 0 ok / 1 missing / 4 user declined install
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
from _lazy import ensure_pkg  # noqa: E402


def extract_slides(path: Path) -> list[dict]:
    ensure_pkg("python-pptx", import_name="pptx")
    from pptx import Presentation  # type: ignore[import-not-found]

    prs = Presentation(str(path))
    out: list[dict] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        body_parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
            if not text:
                continue
            # Heuristic: shape with `title` placeholder type or smallest top-y becomes title
            try:
                ph = shape.placeholder_format
            except Exception:
                ph = None
            if ph is not None and ph.idx == 0 and not title:
                title = text
            else:
                body_parts.append(text)
        if not title and body_parts:
            # First text shape as title fallback
            title = body_parts.pop(0).split("\n")[0]
        notes = ""
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            notes = ""
        body = "\n\n".join(body_parts).strip()
        out.append({
            "source": str(path),
            "slide_index": i,
            "title": title[:200],
            "body": body,
            "notes": notes,
            "word_count": len(body.split()) + len(notes.split()),
        })
    return out


def write(records: list[dict], dest: Path | str) -> None:
    if dest == "-":
        for r in records:
            print(json.dumps(r, ensure_ascii=False))
        return
    p = Path(dest)
    p.parent.mkdir(parents=True, exist_ok=True)
    with p.open("w", encoding="utf-8") as f:
        for r in records:
            f.write(json.dumps(r, ensure_ascii=False) + "\n")
    print(f"pptx_to_chunks: wrote {len(records)} slides → {p}", file=sys.stderr)


def main(argv: list[str] | None = None) -> int:
    p = argparse.ArgumentParser(description="Extract a PPTX deck as JSONL chunks.")
    p.add_argument("input", type=Path)
    p.add_argument("output", nargs="?", default=None)
    args = p.parse_args(argv)

    if not args.input.exists():
        print(f"pptx_to_chunks: input not found: {args.input}", file=sys.stderr)
        return 1

    slides = extract_slides(args.input)
    print(f"pptx_to_chunks: {len(slides)} slides", file=sys.stderr)
    out = args.output
    if out is None:
        out = args.input.with_suffix(args.input.suffix + ".chunks.jsonl")
    write(slides, out)
    return 0


if __name__ == "__main__":
    sys.exit(main())
